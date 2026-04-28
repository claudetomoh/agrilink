"""
AgriLink Ghana — Sprint 3 Final Presentation  →  PowerPoint
CS 415 Software Engineering | April 2026
15 slides: complete platform walkthrough, full Jira board, test results, docs
"""

import os, io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

BASE = os.path.dirname(__file__)

# ── Palette (same as build_pptx.py) ──────────────────────────────────────────
C_PRIMARY   = RGBColor(0x01, 0x2D, 0x1D)
C_PRIMARY_C = RGBColor(0x1B, 0x43, 0x32)
C_ACCENT    = RGBColor(0xFF, 0xC6, 0x41)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_SURFACE   = RGBColor(0xF4, 0xFA, 0xF5)
C_SURFACE_L = RGBColor(0xEC, 0xF3, 0xEC)
C_TEXT      = RGBColor(0x1A, 0x1C, 0x1A)
C_MUTED     = RGBColor(0x4A, 0x55, 0x68)
C_GRAY      = RGBColor(0x9C, 0xA3, 0xAF)
C_TEXT_MUTED= RGBColor(0x4A, 0x55, 0x68)
C_BLUE      = RGBColor(0x1E, 0x40, 0xAF)
C_GOLD      = RGBColor(0x92, 0x6B, 0x00)
C_RED       = RGBColor(0x7E, 0x1D, 0x1D)
C_DARK      = RGBColor(0x06, 0x18, 0x10)
C_MID       = RGBColor(0x0D, 0x20, 0x18)
C_BADGE     = RGBColor(0x0A, 0x1F, 0x14)
C_DIM       = RGBColor(0xAA, 0xBB, 0xAA)
C_DIMMER    = RGBColor(0xCC, 0xDD, 0xCC)

# Extra accent colours
C_GREEN_BADGE = RGBColor(0xDC, 0xFC, 0xE7)
C_GREEN_FG    = RGBColor(0x16, 0x65, 0x34)
C_AMBER_BADGE = RGBColor(0xFE, 0xF3, 0xC7)
C_AMBER_FG    = RGBColor(0x92, 0x40, 0x0E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Drawing helpers ───────────────────────────────────────────────────────────

def bg_solid(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def bg_gradient(slide, c1=C_PRIMARY, c2=C_PRIMARY_C):
    bg_solid(slide, c1)
    shp = slide.shapes.add_shape(1, 0, 0, W, H)
    shp.fill.gradient(); shp.fill.gradient_angle = 135
    shp.fill.gradient_stops[0].position = 0.0; shp.fill.gradient_stops[0].color.rgb = c1
    shp.fill.gradient_stops[1].position = 1.0; shp.fill.gradient_stops[1].color.rgb = c2
    shp.line.fill.background()
    shp._element.getparent().remove(shp._element)
    slide.shapes._spTree.insert(2, shp._element)

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else:    s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h,
       fn="Manrope", fs=Pt(12), bold=False, italic=False,
       color=C_TEXT, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.name = fn; r.font.size = fs; r.font.bold = bold
    r.font.color.rgb = color
    if italic: r.font.italic = True
    return box

def label(slide, text, l, t, dark=False):
    c = C_PRIMARY_C if dark else C_ACCENT
    tb(slide, text, l, t, Inches(8), Inches(0.3),
       fn="Inter", fs=Pt(9), bold=True, color=c)

def heading(slide, text, l, t, color=C_PRIMARY, size=Pt(36)):
    tb(slide, text, l, t, Inches(12), Inches(0.9),
       fs=size, bold=True, color=color)

def embed_screen(slide, path, l, t, w, h, crop_top_px=500):
    """Crop to top crop_top_px of image then embed, letterboxed in the box."""
    img = Image.open(path)
    iw, ih = img.size
    crop_h = min(crop_top_px, ih)
    img = img.crop((0, 0, iw, crop_h))
    target_w = int(w / 914400 * 96)
    target_h = int(h / 914400 * 96)
    ar = iw / crop_h
    if target_w / target_h > ar:
        new_h = target_h; new_w = int(ar * new_h)
    else:
        new_w = target_w; new_h = int(new_w / ar)
    img = img.resize((max(1, new_w), max(1, new_h)), Image.LANCZOS)
    canvas = Image.new("RGB", (target_w, target_h), (6, 24, 16))
    px = (target_w - new_w) // 2; py = (target_h - new_h) // 2
    canvas.paste(img, (px, py))
    buf = io.BytesIO()
    canvas.save(buf, "PNG"); buf.seek(0)
    slide.shapes.add_picture(buf, l, t, w, h)

def embed_img(slide, path, l, t, w, h=None):
    img = Image.open(path)
    iw, ih = img.size
    if h is None:
        h = int(w * ih / iw)
    slide.shapes.add_picture(path, l, t, w, h)

# Screen paths
S = lambda folder: os.path.join(BASE, folder, "screen.png")
SCREENS = {
    "homepage":     S("homepage"),
    "farmer":       S("farmer_dashboard"),
    "buyer":        S("buyer_marketplace"),
    "dashboard":    S("dashboard"),
    "analytics":    S("analytics"),
    "onboarding":   S("onboarding"),
    "logistics":    S("logistics_dashboard"),
    "orders":       S("orders_tracking"),
    "matching":     S("supply_demand_matching"),
    "gh_home":      S("agrilink_ghana_homepage"),
    "mkt_analytic": S("marketplace_analytics"),
}

DIAGRAM_UC   = os.path.join(BASE, "agrilink_usecase.png")
DIAGRAM_ARCH = os.path.join(BASE, "agrilink_arch.png")

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_gradient(s)

embed_screen(s, SCREENS["gh_home"], Inches(8.0), Inches(0), Inches(5.33), H, crop_top_px=700)
rect(s, Inches(7.6), Inches(0), Inches(0.8), H, fill=C_PRIMARY)

rect(s, Inches(0.6), Inches(0.9), Inches(6.4), Inches(0.44), fill=C_BADGE)
tb(s, "CS 415 Software Engineering  ·  Final Presentation  ·  April 2026",
   Inches(0.6), Inches(0.92), Inches(6.4), Inches(0.38),
   fn="Inter", fs=Pt(9), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

tb(s, "AgriLink\nGhana", Inches(0.5), Inches(1.45), Inches(7.3), Inches(2.4),
   fs=Pt(64), bold=True, color=C_WHITE)

tb(s, "Ghana's Agricultural Marketplace — Live & Complete",
   Inches(0.6), Inches(4.0), Inches(7.0), Inches(0.55),
   fn="Inter", fs=Pt(13), color=C_DIM)

pills = ["✅ All 3 Sprints Done", "🌾 4 User Roles", "🏆 34 Stories Shipped", "🌐 Live Deployed"]
for i, p in enumerate(pills):
    px = Inches(0.6) + i * Inches(1.68)
    rect(s, px, Inches(5.0), Inches(1.6), Inches(0.44), fill=C_DARK)
    tb(s, p, px, Inches(5.0), Inches(1.6), Inches(0.44),
       fn="Inter", fs=Pt(9), color=C_DIMMER, align=PP_ALIGN.CENTER)

tb(s, "Tomoh Ikfingeh  ·  claudetomoh/agrilink",
   Inches(0.6), Inches(5.65), Inches(7.0), Inches(0.35),
   fn="Inter", fs=Pt(9), color=C_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 2 — AGENDA
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_solid(s, C_SURFACE)
label(s, "TODAY'S AGENDA", Inches(0.6), Inches(0.3), dark=True)
heading(s, "Final Presentation Outline", Inches(0.6), Inches(0.6))

agenda = [
    ("01", "Platform Summary"),
    ("02", "Sprint 3 Jira Board"),
    ("03", "Matching Engine"),
    ("04", "Architecture"),
    ("05", "Test Results"),
    ("06", "Engineering Docs"),
    ("07", "Live Demo"),
    ("08", "Velocity & Retrospective"),
    ("09", "Roadmap & Q&A"),
]
cw, ch, cgx, cgy = Inches(2.6), Inches(1.8), Inches(0.22), Inches(0.22)
sx, sy = Inches(0.55), Inches(1.7)
for i, (num, name) in enumerate(agenda):
    cx = sx + (i % 3) * (cw + cgx)
    cy = sy + (i // 3) * (ch + cgy)
    rect(s, cx, cy, cw, ch, fill=C_WHITE)
    tb(s, num, cx+Inches(0.18), cy+Inches(0.15), Inches(1), Inches(0.6),
       fs=Pt(36), bold=True, color=RGBColor(0x3D, 0x8C, 0x5E))
    tb(s, name, cx+Inches(0.18), cy+Inches(0.9), cw-Inches(0.36), Inches(0.45),
       fs=Pt(11), bold=True, color=C_PRIMARY)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 3 — PLATFORM SUMMARY
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_gradient(s)

label(s, "PLATFORM SUMMARY", Inches(0.55), Inches(0.35))
tb(s, "AgriLink Ghana — Complete",
   Inches(0.55), Inches(0.65), Inches(8.5), Inches(1.0),
   fs=Pt(38), bold=True, color=C_WHITE)

stats = [
    ("34", "User Stories Shipped"),
    ("83", "Story Points Delivered"),
    ("43", "Tests Passing"),
    ("12", "Produce Listings"),
    ("7",  "Registered Users"),
    ("6",  "Orders Placed"),
]
sw = Inches(4.1); sg = Inches(0.22); sx0 = Inches(0.5)
sh = Inches(1.8)
for i, (val, lab) in enumerate(stats):
    col = i % 3; row = i // 2
    bx = sx0 + col * (sw + sg)
    by = Inches(1.85) + row * (sh + Inches(0.18))
    if i < 3:
        bx = sx0 + col * (sw + sg)
    else:
        bx = sx0 + (col) * (sw + sg)
    rect(s, bx, by, sw, sh, fill=C_DARK)
    rect(s, bx, by, sw, Inches(0.06), fill=C_ACCENT)
    tb(s, val, bx, by+Inches(0.18), sw, Inches(0.95),
       fs=Pt(52), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(s, lab, bx, by+Inches(1.18), sw, Inches(0.42),
       fn="Inter", fs=Pt(10), color=C_DIM, align=PP_ALIGN.CENTER)

roles = [
    ("🌾 Farmer", "Kofi Boateng", "kofi.boateng@agrilink.gh", "Pass@1234"),
    ("🛒 Buyer",   "Kwame Mensah", "kwame.mensah@agrilink.gh",  "Pass@1234"),
    ("🚛 Transport","Kojo Logistics","kojo.logistics@agrilink.gh","Pass@1234"),
    ("🛡️ Admin",   "Admin",        "admin@agrilink.gh",          "Admin@1234"),
]
rw = Inches(3.0); rg = Inches(0.22); rx0 = Inches(0.5)
ry = Inches(6.2)
rect(s, rx0, ry, Inches(12.33), Inches(1.15), fill=C_BADGE)
tb(s, "Demo Accounts  —  Live Server", rx0+Inches(0.2), ry+Inches(0.05),
   Inches(4), Inches(0.32), fn="Inter", fs=Pt(9), bold=True, color=C_ACCENT)
for i, (role, name, email, pwd) in enumerate(roles):
    rx = rx0 + Inches(0.2) + i * (rw + rg)
    tb(s, f"{role}: {email} / {pwd}",
       rx, ry+Inches(0.42), rw+Inches(0.5), Inches(0.28),
       fn="Courier New", fs=Pt(7.5), color=C_DIMMER)
    tb(s, name, rx, ry+Inches(0.72), rw, Inches(0.28),
       fn="Inter", fs=Pt(8), color=C_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 4 — JIRA SPRINT 3 BOARD (All Done)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_solid(s, C_SURFACE)
label(s, "JIRA SPRINT BOARD", Inches(0.6), Inches(0.3), dark=True)
heading(s, "Sprint 3 — All Items Complete", Inches(0.6), Inches(0.6), size=Pt(28))

# Sprint 1&2 done items (moved to final done column) + Sprint 3 items
ALL_DONE = [
    # Sprint 1 & 2
    ("Auth System (Register/Login)",        "Infra"),
    ("Farmer Produce Listings CRUD",        "Farmer"),
    ("Buyer Marketplace + Filters",         "Buyer"),
    ("Price Bidding System",               "Buyer"),
    ("Real-time Notification System",       "Infra"),
    ("Transport Job Board",                "Transport"),
    ("Order Tracking (Full Lifecycle)",     "Buyer"),
    ("Admin User Management",              "Admin"),
    ("Review & Rating System",             "Buyer"),
    # Sprint 3
    ("AI Demand Matching Refinement",       "Buyer"),
    ("Admin Analytics Dashboard",          "Admin"),
    ("Delivery Timeline View",             "Transport"),
    ("Smart Matching Score UI",            "Buyer"),
    ("Security Hardening (OWASP)",         "Infra"),
    ("Full Test Suite — 43 Tests Pass",    "Infra"),
    ("Engineering Documentation Suite",    "Infra"),
]

tag_c = {
    "Farmer":    (RGBColor(0xDC,0xFC,0xE7), RGBColor(0x16,0x65,0x34)),
    "Buyer":     (RGBColor(0xDB,0xEA,0xFE), RGBColor(0x1E,0x40,0xAF)),
    "Admin":     (RGBColor(0xFC,0xE7,0xF3), RGBColor(0x9D,0x17,0x4D)),
    "Transport": (RGBColor(0xFE,0xF3,0xC7), RGBColor(0x92,0x40,0x0E)),
    "Infra":     (RGBColor(0xED,0xE9,0xFE), RGBColor(0x5B,0x21,0xB6)),
}

# Render in 2 columns (done left large, deferred right)
DEFERRED = [
    ("SMS Notification Gateway",           "Infra"),
    ("Ghana Card OCR Verification",        "Farmer"),
    ("Mobile Money (MoMo) Payments",       "Infra"),
    ("PWA Mobile Wrapper",                 "Infra"),
    ("Multi-language: Twi / Ga",           "Infra"),
    ("Dispute Resolution Workflow",        "Admin"),
    ("API Layer for Integrations",         "Infra"),
]

CARD_H = Inches(0.52)
CARD_G = Inches(0.05)

# DONE column — 3 sub-columns
rect(s, Inches(0.4), Inches(1.65), Inches(8.8), Inches(0.42), fill=C_GREEN_BADGE)
tb(s, f"Done ✓  ({len(ALL_DONE)})", Inches(0.4), Inches(1.65), Inches(8.8), Inches(0.42),
   fs=Pt(11), bold=True, color=C_GREEN_FG, align=PP_ALIGN.CENTER)

col_w = Inches(2.82); col_g = Inches(0.12); col_x0 = Inches(0.42)
items_per_col = 6
for ii, (title, tag) in enumerate(ALL_DONE):
    col = ii // items_per_col
    row = ii % items_per_col
    if col > 2: break
    cx = col_x0 + col * (col_w + col_g)
    cy = Inches(2.15) + row * (CARD_H + CARD_G)
    rect(s, cx, cy, col_w, CARD_H, fill=C_WHITE)
    rect(s, cx, cy, col_w, CARD_H, line=C_GREEN_BADGE, lw=Pt(0.5))
    tb(s, "✓", cx+Inches(0.08), cy+Inches(0.11), Inches(0.25), Inches(0.3),
       fn="Inter", fs=Pt(9), bold=True, color=C_GREEN_FG)
    tb(s, title, cx+Inches(0.35), cy+Inches(0.07), col_w-Inches(1.2), Inches(0.38),
       fs=Pt(7.5), bold=True, color=C_PRIMARY)
    if tag in tag_c:
        tbg2, tfg2 = tag_c[tag]
        rect(s, cx+col_w-Inches(0.92), cy+Inches(0.13), Inches(0.84), Inches(0.24), fill=tbg2)
        tb(s, tag, cx+col_w-Inches(0.92), cy+Inches(0.13), Inches(0.84), Inches(0.24),
           fn="Inter", fs=Pt(6.5), bold=True, color=tfg2, align=PP_ALIGN.CENTER)

# DEFERRED column
rect(s, Inches(9.4), Inches(1.65), Inches(3.5), Inches(0.42), fill=C_AMBER_BADGE)
tb(s, f"Sprint 4 Backlog  ({len(DEFERRED)})", Inches(9.4), Inches(1.65), Inches(3.5), Inches(0.42),
   fs=Pt(11), bold=True, color=C_AMBER_FG, align=PP_ALIGN.CENTER)
for ii, (title, tag) in enumerate(DEFERRED):
    cy = Inches(2.15) + ii * (CARD_H + CARD_G)
    rect(s, Inches(9.4), cy, Inches(3.5), CARD_H, fill=C_WHITE)
    rect(s, Inches(9.4), cy, Inches(3.5), CARD_H, line=C_AMBER_BADGE, lw=Pt(0.5))
    tb(s, title, Inches(9.52), cy+Inches(0.07), Inches(2.5), Inches(0.38),
       fs=Pt(8), bold=True, color=C_MUTED)
    if tag in tag_c:
        tbg2, tfg2 = tag_c[tag]
        rect(s, Inches(12.0), cy+Inches(0.13), Inches(0.84), Inches(0.24), fill=tbg2)
        tb(s, tag, Inches(12.0), cy+Inches(0.13), Inches(0.84), Inches(0.24),
           fn="Inter", fs=Pt(6.5), bold=True, color=tfg2, align=PP_ALIGN.CENTER)

# Velocity strip at bottom
vy = Inches(6.85)
rect(s, Inches(0.4), vy, Inches(12.5), Inches(0.5), fill=C_BADGE)
tb(s, "Sprint Velocity:   Sprint 1: 22 pts   ·   Sprint 2: 27 pts   ·   Sprint 3: 34 pts   →   Total: 83 story points delivered",
   Inches(0.6), vy+Inches(0.08), Inches(12.1), Inches(0.35),
   fn="Inter", fs=Pt(10), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 5 — MATCHING ENGINE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_gradient(s)

label(s, "AI SUPPLY–DEMAND MATCHING", Inches(0.55), Inches(0.35))
tb(s, "Smart Score-Based\nMatching Engine",
   Inches(0.55), Inches(0.65), Inches(6.5), Inches(1.6),
   fs=Pt(32), bold=True, color=C_WHITE)

criteria = [
    ("+40", "Same Region",     "Produce and buyer in same Ghana region"),
    ("+30", "Category Match",  "Produce category matches buyer request"),
    ("+20", "Quantity Met",    "Available qty ≥ buyer requested amount"),
    ("+10", "Budget Match",    "Price per unit ≤ buyer budget threshold"),
]
for i, (score, crit, desc) in enumerate(criteria):
    cy = Inches(0.45) + i * Inches(1.5)
    rect(s, Inches(6.8), cy, Inches(0.06), Inches(1.3), fill=C_ACCENT)
    rect(s, Inches(6.9), cy, Inches(6.1), Inches(1.3), fill=C_DARK)
    tb(s, score, Inches(7.0), cy+Inches(0.1), Inches(0.9), Inches(0.5),
       fn="Inter", fs=Pt(26), bold=True, color=C_ACCENT)
    tb(s, crit, Inches(7.95), cy+Inches(0.1), Inches(4.9), Inches(0.38),
       fn="Inter", fs=Pt(13), bold=True, color=C_WHITE)
    tb(s, desc, Inches(7.95), cy+Inches(0.52), Inches(4.9), Inches(0.65),
       fn="Inter", fs=Pt(10.5), color=C_DIM)

# Code block
rect(s, Inches(0.5), Inches(6.35), Inches(6.3), Inches(1.05), fill=RGBColor(0x0D, 0x11, 0x17))
tb(s, "$score = 0;\nif ($region === $buyerRegion)   { $score += 40; }\nif ($category === $requested)     { $score += 30; }\nif ($qty >= $buyerQty)            { $score += 20; }\nif ($price <= $buyerBudget)       { $score += 10; }",
   Inches(0.65), Inches(6.4), Inches(6.0), Inches(0.95),
   fn="Courier New", fs=Pt(8.5), color=RGBColor(0xE6, 0xED, 0xF3))

tb(s, "Max score: 100 pts  ·  Displayed as match % on each listing card  ·  Top-scored results shown first",
   Inches(0.5), Inches(7.2), Inches(6.3), Inches(0.25),
   fn="Inter", fs=Pt(8), color=C_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 6 — ARCHITECTURE DIAGRAM
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_solid(s, C_SURFACE)
label(s, "SYSTEM ARCHITECTURE — MVC LAYERS", Inches(0.4), Inches(0.1), dark=True)
embed_img(s, DIAGRAM_ARCH, Inches(0.3), Inches(0.35), Inches(12.73))

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 7 — TEST RESULTS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_solid(s, C_SURFACE)
label(s, "QUALITY ASSURANCE", Inches(0.6), Inches(0.3), dark=True)
heading(s, "Test Results — 43 / 43 Passing", Inches(0.6), Inches(0.6), size=Pt(28))

# Big pass badge
rect(s, Inches(0.5), Inches(1.55), Inches(3.8), Inches(2.6), fill=C_GREEN_BADGE)
tb(s, "43", Inches(0.5), Inches(1.65), Inches(3.8), Inches(1.4),
   fs=Pt(80), bold=True, color=C_GREEN_FG, align=PP_ALIGN.CENTER)
tb(s, "Tests Passing", Inches(0.5), Inches(3.0), Inches(3.8), Inches(0.5),
   fn="Inter", fs=Pt(14), bold=True, color=C_GREEN_FG, align=PP_ALIGN.CENTER)
tb(s, "0 failures", Inches(0.5), Inches(3.5), Inches(3.8), Inches(0.35),
   fn="Inter", fs=Pt(11), color=C_MUTED, align=PP_ALIGN.CENTER)
rect(s, Inches(0.5), Inches(4.05), Inches(3.8), Inches(0.06), fill=C_GREEN_FG)
tb(s, "100% pass rate", Inches(0.5), Inches(4.2), Inches(3.8), Inches(0.35),
   fn="Inter", fs=Pt(10), bold=True, color=C_GREEN_FG, align=PP_ALIGN.CENTER)

# Test breakdown table
test_groups = [
    ("Authentication",      8, ["Login/Logout flow", "CSRF token validation", "Role-based redirects", "bcrypt verify"]),
    ("Farmer Module",       9, ["Listing CRUD operations", "Image upload validation", "Order accept/reject", "Profile edit"]),
    ("Buyer Module",        8, ["Marketplace filters", "Order placement", "Bid counter-offer", "Review submission"]),
    ("Transport Module",    6, ["Job board listing", "Delivery status update", "Timeline view"]),
    ("Analytics & Admin",   6, ["Dashboard metrics", "User toggle active", "Order/delivery views"]),
    ("Security (OWASP)",    6, ["SQL injection (PDO)", "XSS output escaping", "Auth guards", "File upload rules"]),
]
tw = Inches(1.53); tg = Inches(0.15); tx0 = Inches(4.6)
for ti, (name, count, checks) in enumerate(test_groups):
    col = ti % 3; row = ti // 3
    tx = tx0 + col * (tw * 2.85 + tg * 0.5)
    ty = Inches(1.55) + row * Inches(2.75)
    bw = tw * 2.85
    rect(s, tx, ty, bw, Inches(0.4), fill=C_PRIMARY_C)
    tb(s, f"{name}  ({count})", tx+Inches(0.1), ty+Inches(0.02), bw-Inches(0.2), Inches(0.36),
       fs=Pt(9.5), bold=True, color=C_WHITE)
    rect(s, tx, ty+Inches(0.4), bw, Inches(2.15), fill=C_WHITE)
    for ci, chk in enumerate(checks):
        tb(s, "✓  " + chk, tx+Inches(0.12), ty+Inches(0.5)+ci*Inches(0.5),
           bw-Inches(0.24), Inches(0.42),
           fn="Inter", fs=Pt(8.5), color=C_MUTED)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 8 — ENGINEERING DOCUMENTATION
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_solid(s, C_SURFACE)
label(s, "PROJECT DOCUMENTATION", Inches(0.6), Inches(0.3), dark=True)
heading(s, "Engineering Documentation Suite", Inches(0.6), Inches(0.6), size=Pt(28))

docs = [
    ("📋", "README.md",             "Project overview, setup guide, architecture map, credential table, all 7 engineering docs index"),
    ("📄", "PROJECT_PROPOSAL.md",   "Problem statement, objectives & success metrics, tech stack, 7-week timeline, team Scrum roles (Week 3)"),
    ("🏃", "SPRINT.md",             "3-sprint log — Jira board state, daily standups, velocity chart (22→27→34 pts), retrospectives, DoD"),
    ("📑", "FINAL_REPORT.md",       "11-section report: overview, methodology, architecture, OWASP table, defect log, 43-test results, learnings"),
    ("🗄️", "schema.sql",            "7-table MySQL schema with FK constraints, ENUMs, indexes — production-ready"),
    ("🌱", "seed.sql",              "Ghana-context demo data: Ghanaian names, produce, 16 regions, test accounts"),
    ("📐", "DESIGN.md",             "Design system — Verdant Archive palette, typography, component library, spacing tokens"),
]

dw = Inches(12.0); dx0 = Inches(0.6)
for di, (ico, name, desc) in enumerate(docs):
    dy = Inches(1.65) + di * Inches(0.76)
    alt = (di % 2 == 0)
    rect(s, dx0, dy, dw, Inches(0.72), fill=C_WHITE if alt else C_SURFACE_L)
    tb(s, ico, dx0+Inches(0.12), dy+Inches(0.16), Inches(0.38), Inches(0.38),
       fs=Pt(14))
    tb(s, name, dx0+Inches(0.6), dy+Inches(0.05), Inches(2.8), Inches(0.36),
       fn="Courier New", fs=Pt(10.5), bold=True, color=C_PRIMARY)
    tb(s, desc, dx0+Inches(3.5), dy+Inches(0.05), Inches(8.4), Inches(0.62),
       fn="Inter", fs=Pt(9.5), color=C_MUTED)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 9 — LIVE APPLICATION SCREENS (Homepage + Buyer Marketplace)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_gradient(s, C_MID, C_PRIMARY_C)
label(s, "LIVE APPLICATION — PORTAL SCREENS", Inches(0.5), Inches(0.2))
heading(s, "Homepage · Buyer Marketplace · AI Matching", Inches(0.5), Inches(0.47),
        color=C_WHITE, size=Pt(22))

sw = Inches(4.1); sh = Inches(5.8); sg = Inches(0.22); sx0 = Inches(0.5)
set_a = [(SCREENS["gh_home"],"AgriLink Homepage",600),(SCREENS["buyer"],"Buyer Marketplace",600),(SCREENS["matching"],"AI Supply-Demand Matching",600)]
for i, (path, caption, crop) in enumerate(set_a):
    sx = sx0 + i*(sw+sg)
    rect(s, sx, Inches(1.15), sw, Inches(0.32), fill=RGBColor(0x20,0x30,0x28))
    tb(s, "● ● ●", sx+Inches(0.1), Inches(1.15), Inches(0.7), Inches(0.32),
       fn="Inter", fs=Pt(8), color=RGBColor(0x66,0x88,0x77))
    embed_screen(s, path, sx, Inches(1.47), sw, sh-Inches(0.32), crop_top_px=crop)
    tb(s, caption, sx, Inches(7.05), sw, Inches(0.35),
       fn="Inter", fs=Pt(10), bold=True, color=C_DIMMER, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 10 — LIVE APPLICATION SCREENS (Farmer + Admin + Transport)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_gradient(s, C_MID, C_PRIMARY_C)
label(s, "LIVE APPLICATION — PORTAL SCREENS", Inches(0.5), Inches(0.2))
heading(s, "Farmer Dashboard · Admin Dashboard · Transport Dashboard", Inches(0.5), Inches(0.47),
        color=C_WHITE, size=Pt(22))

set_b = [(SCREENS["farmer"],"Farmer Dashboard",600),(SCREENS["dashboard"],"Admin Dashboard",600),(SCREENS["logistics"],"Transport Dashboard",600)]
for i, (path, caption, crop) in enumerate(set_b):
    sx = sx0 + i*(sw+sg)
    rect(s, sx, Inches(1.15), sw, Inches(0.32), fill=RGBColor(0x20,0x30,0x28))
    tb(s, "● ● ●", sx+Inches(0.1), Inches(1.15), Inches(0.7), Inches(0.32),
       fn="Inter", fs=Pt(8), color=RGBColor(0x66,0x88,0x77))
    embed_screen(s, path, sx, Inches(1.47), sw, sh-Inches(0.32), crop_top_px=crop)
    tb(s, caption, sx, Inches(7.05), sw, Inches(0.35),
       fn="Inter", fs=Pt(10), bold=True, color=C_DIMMER, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 11 — VELOCITY & RETROSPECTIVE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_solid(s, C_SURFACE)
label(s, "AGILE PROCESS — RETROSPECTIVE", Inches(0.6), Inches(0.3), dark=True)
heading(s, "Sprint Velocity & Key Learnings", Inches(0.6), Inches(0.6), size=Pt(28))

# Velocity bars
sprints = [("Sprint 1", 22, 27), ("Sprint 2", 27, 27), ("Sprint 3", 34, 27)]
bar_max_w = Inches(4.8)
for i, (sname, pts, _) in enumerate(sprints):
    by = Inches(1.7) + i * Inches(1.3)
    bar_w = bar_max_w * (pts / 40)
    rect(s, Inches(0.6), by, Inches(1.5), Inches(0.9), fill=C_PRIMARY_C)
    tb(s, sname, Inches(0.6), by+Inches(0.24), Inches(1.5), Inches(0.42),
       fn="Inter", fs=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.3), by+Inches(0.2), bar_w, Inches(0.5), fill=C_ACCENT)
    tb(s, f"{pts} story points", Inches(2.35), by+Inches(0.22), bar_w, Inches(0.42),
       fn="Inter", fs=Pt(11), bold=True, color=C_PRIMARY)

rect(s, Inches(0.6), Inches(5.65), Inches(5.0), Inches(0.06), fill=C_ACCENT)
tb(s, "Velocity trend: +23% Sprint 1→2, +26% Sprint 2→3", Inches(0.6), Inches(5.78),
   Inches(5.0), Inches(0.35), fn="Inter", fs=Pt(10), color=C_MUTED)

# Retrospective highlights
retros = [
    ("Sprint 1", [
        "✅ Went well: Auth, CSRF, listings — clean separation of concerns",
        "🔧 Improved: PHP lint gate + deploy script replaced manual FTP",
        "📌 Action: Define DoD before sprint start; review design tokens",
    ]),
    ("Sprint 2", [
        "✅ Went well: Order lifecycle, bidding, transport job board",
        "🔧 Improved: CSS contrast audit fixed 3 accessibility failures",
        "📌 Action: Add unit tests from Sprint 3 start — not at the end",
    ]),
    ("Sprint 3", [
        "✅ Went well: Matching engine, analytics, full 43-test suite",
        "🔧 Improved: Documentation written alongside code, not after",
        "📌 Carry forward: SMS gateway + OCR to Sprint 4 with spec ready",
    ]),
]
rcw = Inches(3.9); rcx0 = Inches(6.6); rcy = Inches(1.55)
for i, (sname, bullets) in enumerate(retros):
    rx = rcx0 + i * (rcw + Inches(0.18))
    rect(s, rx, rcy, rcw, Inches(0.4), fill=C_PRIMARY_C)
    tb(s, sname, rx+Inches(0.15), rcy+Inches(0.04), rcw, Inches(0.32),
       fn="Inter", fs=Pt(11), bold=True, color=C_WHITE)
    rect(s, rx, rcy+Inches(0.4), rcw, Inches(4.5), fill=C_WHITE)
    for bi, blt in enumerate(bullets):
        tb(s, blt, rx+Inches(0.15), rcy+Inches(0.5)+bi*Inches(1.3),
           rcw-Inches(0.3), Inches(1.2),
           fn="Inter", fs=Pt(9), color=C_MUTED)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 12 — SECURITY & CODE QUALITY (OWASP)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_solid(s, C_SURFACE)
label(s, "SECURITY — OWASP TOP 10", Inches(0.6), Inches(0.3), dark=True)
heading(s, "Security Implementation", Inches(0.6), Inches(0.6), size=Pt(28))

owasp = [
    ("A01 Broken Access Control",   "Auth::requireRole() guard on every controller action"),
    ("A02 Cryptographic Failures",  "password_hash(BCRYPT) — no plain text ever stored"),
    ("A03 Injection",               "PDO prepared statements on all DB queries"),
    ("A04 Insecure Design",         "Dedicated Auth, Session, CSRF token classes"),
    ("A05 Security Misconfiguration","DB credentials in config.php excluded from version control"),
    ("A07 Auth Failures",           "Session::verifyCsrf() on every POST; session_regenerate_id()"),
    ("A08 Integrity Failures",      "File upload: MIME + ext + size checks; store outside webroot"),
    ("A09 Logging Failures",        "PHP error_log; future: audit log table for admin actions"),
]
ow = Inches(12.0); ox0 = Inches(0.6)
for oi, (threat, mitigation) in enumerate(owasp):
    oy = Inches(1.65) + oi * Inches(0.68)
    alt = (oi % 2 == 0)
    rect(s, ox0, oy, ow, Inches(0.64), fill=C_WHITE if alt else C_SURFACE_L)
    rect(s, ox0, oy, Inches(0.06), Inches(0.64), fill=C_GREEN_FG)
    tb(s, threat, ox0+Inches(0.18), oy+Inches(0.04), Inches(4.0), Inches(0.56),
       fn="Inter", fs=Pt(9.5), bold=True, color=C_PRIMARY)
    tb(s, mitigation, ox0+Inches(4.3), oy+Inches(0.04), Inches(7.6), Inches(0.56),
       fn="Inter", fs=Pt(9.5), color=C_MUTED)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 13 — LIVE DEMO
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_gradient(s)

embed_screen(s, SCREENS["buyer"], Inches(7.8), Inches(0), Inches(5.53), H, crop_top_px=700)
rect(s, Inches(7.4), Inches(0), Inches(0.7), H, fill=C_PRIMARY)

label(s, "LIVE DEPLOYMENT", Inches(0.6), Inches(0.5))
heading(s, "Application Demo", Inches(0.6), Inches(0.8), color=C_WHITE, size=Pt(34))

rect(s, Inches(0.6), Inches(1.75), Inches(6.5), Inches(0.55), fill=C_DARK)
tb(s, "🌐  169.239.251.102:280/~tomoh.ikfingeh/agrilink/public/",
   Inches(0.6), Inches(1.75), Inches(6.5), Inches(0.55),
   fn="Courier New", fs=Pt(9), color=C_ACCENT, align=PP_ALIGN.CENTER)

accounts = [
    ("🌾 Farmer",    "kofi.boateng@agrilink.gh",   "Pass@1234"),
    ("🛒 Buyer",     "kwame.mensah@agrilink.gh",    "Pass@1234"),
    ("🚛 Transport", "kojo.logistics@agrilink.gh",  "Pass@1234"),
    ("🛡️ Admin",     "admin@agrilink.gh",           "Admin@1234"),
]
aw = Inches(3.0); ag = Inches(0.2)
for ai, (role, email, pwd) in enumerate(accounts):
    col = ai % 2; row = ai // 2
    ax = Inches(0.6) + col*(aw+ag); ay = Inches(2.55) + row * Inches(1.65)
    rect(s, ax, ay, aw, Inches(1.45), fill=C_DARK)
    tb(s, role,  ax+Inches(0.15), ay+Inches(0.1),  aw-Inches(0.3), Inches(0.38),
       fn="Inter", fs=Pt(11), bold=True, color=C_WHITE)
    tb(s, email, ax+Inches(0.15), ay+Inches(0.52), aw-Inches(0.3), Inches(0.32),
       fn="Courier New", fs=Pt(8.5), color=C_ACCENT)
    tb(s, f"pwd: {pwd}", ax+Inches(0.15), ay+Inches(0.88), aw-Inches(0.3), Inches(0.3),
       fn="Courier New", fs=Pt(8.5), color=C_GRAY)

rect(s, Inches(0.6), Inches(6.0), Inches(6.5), Inches(0.48), fill=C_DARK)
tb(s, "📂  github.com/claudetomoh/agrilink",
   Inches(0.6), Inches(6.0), Inches(6.5), Inches(0.48),
   fn="Courier New", fs=Pt(9.5), color=C_ACCENT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 14 — ROADMAP (Sprint 4 & Future)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_solid(s, C_SURFACE)
label(s, "ROADMAP", Inches(0.6), Inches(0.3), dark=True)
heading(s, "Deferred Backlog — Jira Sprint 4 Queue", Inches(0.6), Inches(0.6), size=Pt(28))

next_cols = [
    (C_PRIMARY_C, "Sprint 4 — Trust & Verification", [
        "Ghana Card OCR auto-verification",
        "SMS alerts via Twilio",
        "Email notification system",
        "Dispute resolution workflow",
    ]),
    (C_GOLD,      "Sprint 5 — Growth", [
        "Mobile Money (MoMo) payment integration",
        "Market price index dashboard",
        "Bulk order management",
        "Farmer cooperative groups",
    ]),
    (C_BLUE,      "Future — Scale", [
        "PWA mobile wrapper",
        "Multi-language: Twi / Ga / Hausa",
        "API layer for third-party integrations",
        "ML demand forecasting model",
    ]),
]
ncw = Inches(4.0); ncg = Inches(0.24); ncx0 = Inches(0.5)
for ni, (nc, sprint, items) in enumerate(next_cols):
    nx = ncx0 + ni*(ncw+ncg)
    rect(s, nx, Inches(1.65), ncw, Inches(0.055), fill=nc)
    rect(s, nx, Inches(1.7),  ncw, Inches(5.5),   fill=C_WHITE)
    tb(s, sprint, nx+Inches(0.18), Inches(1.82), ncw-Inches(0.36), Inches(0.5),
       fn="Inter", fs=Pt(10.5), bold=True, color=nc)
    for ii, item in enumerate(items):
        iy = Inches(2.45) + ii * Inches(1.1)
        rect(s, nx+Inches(0.18), iy, Inches(0.06), Inches(0.4), fill=nc)
        tb(s, item, nx+Inches(0.35), iy, ncw-Inches(0.55), Inches(0.4),
           fn="Inter", fs=Pt(10.5), color=C_MUTED)

# Jira note
rect(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.38), fill=C_BADGE)
tb(s, "All Sprint 4 items are tracked as 'To Do' tickets on the Jira sprint board — backlog-ready with acceptance criteria",
   Inches(0.65), Inches(7.04), Inches(12.0), Inches(0.3),
   fn="Inter", fs=Pt(9), color=C_DIM)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 15 — THANK YOU
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_gradient(s)

tb(s, "🌿", Inches(5.9), Inches(0.6), Inches(1.53), Inches(1.1),
   fs=Pt(52), align=PP_ALIGN.CENTER, color=C_WHITE)

rect(s, Inches(3.7), Inches(1.8), Inches(5.93), Inches(0.42), fill=C_BADGE)
tb(s, "AgriLink Ghana  ·  CS 415 Software Engineering  ·  Final Presentation  ·  April 2026",
   Inches(3.7), Inches(1.85), Inches(5.93), Inches(0.35),
   fn="Inter", fs=Pt(9), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

tb(s, "Thank You", Inches(1), Inches(2.4), Inches(11.33), Inches(1.4),
   fs=Pt(72), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

tb(s, "Questions?", Inches(1), Inches(3.8), Inches(11.33), Inches(0.55),
   fn="Inter", fs=Pt(20), color=RGBColor(0x66,0x88,0x77), align=PP_ALIGN.CENTER)

# Key stats strip
stats_row = [("34", "Stories"), ("83", "Pts"), ("43/43", "Tests"), ("3", "Sprints"), ("4", "Portals"), ("Live", "Deployed")]
sw2 = Inches(1.88); sg2 = Inches(0.14); sx2 = Inches(0.7)
sy2 = Inches(4.55)
for i, (val, lbl) in enumerate(stats_row):
    bx = sx2 + i*(sw2+sg2)
    rect(s, bx, sy2, sw2, Inches(0.7), fill=C_DARK)
    tb(s, val, bx, sy2+Inches(0.03), sw2, Inches(0.38),
       fn="Inter", fs=Pt(16), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    tb(s, lbl, bx, sy2+Inches(0.4), sw2, Inches(0.28),
       fn="Inter", fs=Pt(8), color=C_GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(2.5), Inches(5.5), Inches(8.33), Inches(0.5), fill=C_DARK)
tb(s, "🌐  169.239.251.102:280/~tomoh.ikfingeh/agrilink/public/",
   Inches(2.5), Inches(5.5), Inches(8.33), Inches(0.5),
   fn="Courier New", fs=Pt(10), color=C_ACCENT, align=PP_ALIGN.CENTER)

rect(s, Inches(2.5), Inches(6.15), Inches(8.33), Inches(0.5), fill=C_DARK)
tb(s, "📂  github.com/claudetomoh/agrilink",
   Inches(2.5), Inches(6.15), Inches(8.33), Inches(0.5),
   fn="Courier New", fs=Pt(10), color=C_ACCENT, align=PP_ALIGN.CENTER)

tb(s, "Tomoh Ikfingeh  ·  Student  ·  CS 415 Software Engineering",
   Inches(2.5), Inches(6.85), Inches(8.33), Inches(0.32),
   fn="Inter", fs=Pt(9), color=C_GRAY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(BASE, "AgriLink_Final_Presentation.pptx")
prs.save(out)
print(f"✅  Saved: {out}")
print(f"   Slides: 15")
print(f"   Covers: Platform Summary, Jira Sprint 3 (all done), Matching Engine,")
print(f"           Architecture, Test Results (43/43), Engineering Docs,")
print(f"           Live Demo (correct credentials), Velocity, OWASP, Roadmap")
